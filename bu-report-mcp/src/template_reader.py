"""
One-time analysis script: extracts slide structure, colors, and fonts
from templates/utilities_reference.pptx to guide ppt_writer.py redesign.

Usage:
    python src/template_reader.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json


TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "templates" / "utilities_reference.pptx"


def rgb_to_hex(rgb):
    if rgb is None:
        return None
    return f"#{rgb.r:02X}{rgb.g:02X}{rgb.b:02X}"


def extract_run_format(run):
    fmt = {}
    try:
        fmt["bold"] = run.font.bold
        fmt["italic"] = run.font.italic
        fmt["size_pt"] = run.font.size.pt if run.font.size else None
        fmt["name"] = run.font.name
        fmt["color"] = rgb_to_hex(run.font.color.rgb) if run.font.color and run.font.color.type else None
    except Exception:
        pass
    return fmt


def extract_shape_info(shape):
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left_emu": shape.left,
        "top_emu": shape.top,
        "width_emu": shape.width,
        "height_emu": shape.height,
    }

    # Fill color
    try:
        fill = shape.fill
        info["fill_type"] = str(fill.type)
        if fill.type is not None:
            try:
                info["fill_color"] = rgb_to_hex(fill.fore_color.rgb)
            except Exception:
                info["fill_color"] = None
    except Exception:
        pass

    # Text content
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            para_info = {
                "text": para.text,
                "alignment": str(para.alignment),
                "runs": [],
            }
            for run in para.runs:
                para_info["runs"].append({
                    "text": run.text,
                    "format": extract_run_format(run),
                })
            paragraphs.append(para_info)
        info["text_frame"] = paragraphs

    # Table
    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": table.rows.__len__(),
            "cols": table.columns.__len__(),
            "sample_cells": [
                {"row": 0, "col": c, "text": table.cell(0, c).text}
                for c in range(min(3, table.columns.__len__()))
            ],
        }

    # Chart
    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
        }

    return info


def extract_theme_colors(prs):
    """Extract theme colors from the presentation's XML theme."""
    colors = {}
    try:
        theme = prs.core_properties  # not the right path, use slide master
        slide_master = prs.slide_masters[0]
        theme_element = slide_master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        if theme_element is None:
            # Try direct theme from slide master relationships
            for rel in slide_master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel._target
                    clrScheme = theme_part._element.find(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme"
                    )
                    if clrScheme is not None:
                        for child in clrScheme:
                            tag = child.tag.split("}")[-1]
                            srgb = child.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                            )
                            sysClr = child.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr"
                            )
                            if srgb is not None:
                                colors[tag] = f"#{srgb.get('val', '').upper()}"
                            elif sysClr is not None:
                                colors[tag] = f"#{sysClr.get('lastClr', '').upper()}"
    except Exception as e:
        colors["_error"] = str(e)
    return colors


def main():
    if not TEMPLATE_PATH.exists():
        print(f"Template not found: {TEMPLATE_PATH}")
        return

    prs = Presentation(str(TEMPLATE_PATH))

    print("=" * 70)
    print(f"TEMPLATE: {TEMPLATE_PATH.name}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print("=" * 70)

    # Theme colors
    print("\n--- THEME COLORS ---")
    theme_colors = extract_theme_colors(prs)
    for name, val in theme_colors.items():
        print(f"  {name}: {val}")

    # Slide layouts available
    print("\n--- SLIDE LAYOUTS (from Slide Master) ---")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")

    # Each slide
    print("\n--- SLIDES ---")
    for slide_idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        print(f"\n[Slide {slide_idx + 1}] layout='{layout_name}'")
        print(f"  Shapes: {len(slide.shapes)}")

        for shape in slide.shapes:
            info = extract_shape_info(shape)
            left_in = (info.get("left_emu") or 0) / 914400
            top_in = (info.get("top_emu") or 0) / 914400
            w_in = (info.get("width_emu") or 0) / 914400
            h_in = (info.get("height_emu") or 0) / 914400

            print(f"\n  Shape [{info['shape_id']}] '{info['name']}' ({info['shape_type']})")
            print(f"    pos=({left_in:.2f}\", {top_in:.2f}\")  size=({w_in:.2f}\" x {h_in:.2f}\")")

            if "fill_color" in info and info["fill_color"]:
                print(f"    fill={info['fill_color']}")

            if "text_frame" in info:
                for para in info["text_frame"]:
                    text = para["text"].strip()
                    if not text:
                        continue
                    align = para["alignment"]
                    print(f"    text: \"{text[:80]}\"  align={align}")
                    for run in para["runs"]:
                        fmt = run["format"]
                        size = fmt.get("size_pt")
                        font = fmt.get("name")
                        color = fmt.get("color")
                        bold = fmt.get("bold")
                        print(f"      run: \"{run['text'][:60]}\"  font={font} size={size}pt bold={bold} color={color}")

            if "table" in info:
                t = info["table"]
                print(f"    table: {t['rows']}r x {t['cols']}c  sample={t['sample_cells']}")

            if "chart" in info:
                c = info["chart"]
                print(f"    chart: type={c['chart_type']}  legend={c['has_legend']}")

    print("\n" + "=" * 70)
    print("Done.")


if __name__ == "__main__":
    main()
