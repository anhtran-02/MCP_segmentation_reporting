from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Color palette — matches utilities_reference.pptx
BG          = RGBColor(0x0D, 0x2B, 0x3E)   # dark navy background
TEAL        = RGBColor(0x0F, 0xA3, 0xB1)   # left accent bar + highlights
GOLD        = RGBColor(0xF7, 0xB7, 0x31)   # KPI numbers + separators
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE  = RGBColor(0xA8, 0xC8, 0xDC)   # subtitle / secondary text
MUTED       = RGBColor(0x9A, 0xBC, 0xCC)   # KPI labels + secondary labels
KPI_BOX_BG  = RGBColor(0x16, 0x32, 0x48)   # KPI card background
KPI_BOX_BD  = RGBColor(0x2E, 0x5A, 0x7A)   # KPI card border
CARD_BORDER = RGBColor(0xCB, 0xD5, 0xE1)   # agenda card border
RED_HIGH    = RGBColor(0xE7, 0x4C, 0x3C)   # high-priority badge

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.62)

RECT = 1  # MSO rectangle shape type integer


def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(RECT, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _run(tf_or_para, text, size_pt, bold=False, color=WHITE, font=None, align=None, first=False):
    """Add a paragraph+run to a text frame.  first=True uses paragraph[0] instead of adding."""
    if hasattr(tf_or_para, "paragraphs"):
        p = tf_or_para.paragraphs[0] if first else tf_or_para.add_paragraph()
    else:
        p = tf_or_para
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.color.rgb = color
    if font:
        r.font.name = font
    return p


def _inner_chrome(slide, title_text):
    """Header bar + left teal bar + white title — used on all inner slides."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.70), BG)
    _rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, TEAL)
    tb = _tb(slide, Inches(0.35), Inches(0.12), Inches(9.30), Inches(0.50))
    _run(tb.text_frame, title_text, 13, bold=True, color=WHITE, first=True)


def _cover_chrome(slide):
    """Left teal bar + gold top line — used on cover slide only."""
    _rect(slide, Inches(0), Inches(0), Inches(0.14), SLIDE_H, TEAL)
    _rect(slide, Inches(0.14), Inches(0), Inches(9.86), Inches(0.05), GOLD)


def _kpi_box(slide, left, top, width, height, label, value):
    _rect(slide, left, top, width, height, KPI_BOX_BG, KPI_BOX_BD, 0.5)
    _rect(slide, left, top, width, Inches(0.07), TEAL)
    tb_n = _tb(slide, left, top + Inches(0.10), width, Inches(0.70))
    _run(tb_n.text_frame, value, 30, bold=True, color=GOLD, font="Georgia",
         align=PP_ALIGN.CENTER, first=True)
    tb_l = _tb(slide, left, top + Inches(0.80), width, Inches(0.28))
    _run(tb_l.text_frame, label, 11, bold=True, color=MUTED,
         align=PP_ALIGN.CENTER, first=True)


# ─── Slide 1: Cover ──────────────────────────────────────────────────────────

def add_cover_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _cover_chrome(slide)

    kpis = report_data["team_overview"]["kpis"]

    tb1 = _tb(slide, Inches(0.45), Inches(0.60), Inches(9.0), Inches(0.35))
    _run(tb1.text_frame, "ANALYTICS REPORT  ·  DATA TEAM", 9.5, bold=True, color=TEAL, first=True)

    tb2 = _tb(slide, Inches(0.45), Inches(1.00), Inches(9.0), Inches(1.10))
    _run(tb2.text_frame, f"{report_data['bu_name']} Team", 52, bold=True, color=WHITE,
         font="Georgia", first=True)

    tb3 = _tb(slide, Inches(0.45), Inches(2.15), Inches(9.0), Inches(0.55))
    _run(tb3.text_frame, "Segment & Campaign Pattern Analysis", 22, color=LIGHT_BLUE, first=True)

    _rect(slide, Inches(0.45), Inches(2.85), Inches(2.80), Inches(0.05), GOLD)

    tb4 = _tb(slide, Inches(0.45), Inches(3.05), Inches(9.0), Inches(0.38))
    p4 = tb4.text_frame.paragraphs[0]
    r = p4.add_run(); r.text = "Period: "; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = MUTED
    r2 = p4.add_run(); r2.text = report_data["period"]; r2.font.size = Pt(11); r2.font.color.rgb = WHITE

    kpi_xs = [Inches(0.45), Inches(2.83), Inches(5.21), Inches(7.59)]
    kpi_labels = ["Campaigns", "Segments", "Operators", "Avg Cond/Seg"]
    kpi_vals = [kpis["total_campaigns"], kpis["total_segments"],
                kpis["total_operators"], kpis["avg_cond_per_seg"]]
    for x, lbl, val in zip(kpi_xs, kpi_labels, kpi_vals):
        _kpi_box(slide, x, Inches(3.70), Inches(2.18), Inches(1.62), lbl, val)


# ─── Slide 2: Agenda ─────────────────────────────────────────────────────────

def add_agenda_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "AGENDA")

    total_camps = report_data["team_overview"]["kpis"]["total_campaigns"]
    items = [
        ("01", "Methodology & Terminology", "Data sources, join logic, dedup, key definitions"),
        ("02", "Team Overview", f"{total_camps} campaigns, channel breakdown, goal distribution"),
        ("03", "Condition Analysis", "Top conditions, CDP vs BQ, targeting depth"),
        ("04", "Size Pool Distribution", "Audience size by goal × channel, skewness"),
        ("05", "Per-Person Tactics", "Operating profiles & complexity labels"),
        ("06", "Segment Quality Issues", "Zero-condition segments, reuse patterns"),
        ("07", "Recommendations", "Actionable improvements with priority"),
    ]

    col_xs = [Inches(0.30), Inches(5.30)]
    row_counts = [4, 3]
    item_w = Inches(4.60)
    item_h = Inches(0.98)
    gap_y = Inches(1.15)
    start_y = Inches(0.90)

    idx = 0
    for col_x, n_rows in zip(col_xs, row_counts):
        for row_i in range(n_rows):
            num, title, desc = items[idx]; idx += 1
            y = start_y + row_i * gap_y
            _rect(slide, col_x, y, item_w, item_h, WHITE, CARD_BORDER, 0.5)
            _rect(slide, col_x, y, Inches(0.08), item_h, TEAL)
            tb_n = _tb(slide, col_x + Inches(0.10), y, Inches(0.50), item_h)
            _run(tb_n.text_frame, num, 20, bold=True, color=BG, font="Georgia",
                 align=PP_ALIGN.CENTER, first=True)
            tb_t = _tb(slide, col_x + Inches(0.67), y + Inches(0.10), Inches(3.75), Inches(0.35))
            _run(tb_t.text_frame, title, 13, bold=True, color=BG, first=True)
            tb_d = _tb(slide, col_x + Inches(0.67), y + Inches(0.52), Inches(3.75), Inches(0.35))
            _run(tb_d.text_frame, desc, 9, color=RGBColor(0x44, 0x44, 0x44), first=True)


# ─── Slide 3: Methodology ────────────────────────────────────────────────────

def add_methodology_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "01  METHODOLOGY & TERMINOLOGY")

    m = report_data["methodology"]

    # Left: DATA METHODOLOGY section
    left_x = Inches(0.30)
    sec_y = Inches(0.80)
    _rect(slide, left_x, sec_y, Inches(4.50), Inches(0.30), TEAL)
    tb_sec = _tb(slide, left_x + Inches(0.10), sec_y + Inches(0.04), Inches(4.30), Inches(0.22))
    _run(tb_sec.text_frame, "DATA METHODOLOGY", 10, bold=True, color=BG, first=True)

    meth_items = [
        ("Data Sources", m.get("data_sources", "")),
        ("Join Logic", m.get("join_logic", "")),
        ("Dedup Logic", m.get("dedup_logic", "")),
        ("Condition Label", m.get("condition_label_formula", "")),
        ("Key Metrics", m.get("key_metrics", "")),
    ]
    item_y = Inches(1.15)
    for lbl, val in meth_items:
        _rect(slide, left_x, item_y, Inches(4.50), Inches(0.80), KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, left_x, item_y, Inches(0.07), Inches(0.80), TEAL)
        tb_l = _tb(slide, left_x + Inches(0.13), item_y + Inches(0.05), Inches(3.90), Inches(0.24))
        _run(tb_l.text_frame, lbl, 10, bold=True, color=GOLD, first=True)
        tb_v = _tb(slide, left_x + Inches(0.13), item_y + Inches(0.30), Inches(3.90), Inches(0.44))
        tb_v.text_frame.word_wrap = True
        _run(tb_v.text_frame, str(val)[:140], 8.5, color=LIGHT_BLUE, first=True)
        item_y += Inches(0.88)

    # Right: KEY TERMINOLOGY section
    right_x = Inches(5.05)
    _rect(slide, right_x, sec_y, Inches(4.65), Inches(0.30), TEAL)
    tb_sec2 = _tb(slide, right_x + Inches(0.10), sec_y + Inches(0.04), Inches(4.45), Inches(0.22))
    _run(tb_sec2.text_frame, "KEY TERMINOLOGY", 10, bold=True, color=BG, first=True)

    terms = m.get("terminology", [])
    term_y = Inches(1.15)
    for term in terms[:7]:
        k, v = (term.split(":", 1) + [""])[:2]
        _rect(slide, right_x, term_y, Inches(4.65), Inches(0.50), KPI_BOX_BG, KPI_BOX_BD, 0.5)
        tb_k = _tb(slide, right_x + Inches(0.10), term_y + Inches(0.04), Inches(1.40), Inches(0.20))
        _run(tb_k.text_frame, k.strip(), 9.5, bold=True, color=GOLD, first=True)
        tb_v = _tb(slide, right_x + Inches(0.10), term_y + Inches(0.26), Inches(4.45), Inches(0.22))
        tb_v.text_frame.word_wrap = True
        _run(tb_v.text_frame, v.strip()[:120], 8.5, color=LIGHT_BLUE, first=True)
        term_y += Inches(0.58)


# ─── Slide 4: Team Overview ───────────────────────────────────────────────────

def add_team_overview_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "02  TEAM OVERVIEW — Channel & Goal Distribution")

    overview = report_data["team_overview"]
    kpis = overview["kpis"]

    kpi_xs = [Inches(0.30), Inches(2.68), Inches(5.06), Inches(7.44)]
    kpi_labels = ["Campaigns", "Segments", "Operators", "Avg Cond/Seg"]
    kpi_vals = [kpis["total_campaigns"], kpis["total_segments"],
                kpis["total_operators"], kpis["avg_cond_per_seg"]]
    kpi_y = Inches(0.82)
    kpi_h = Inches(1.05)
    kpi_w = Inches(2.18)
    for x, lbl, val in zip(kpi_xs, kpi_labels, kpi_vals):
        _rect(slide, x, kpi_y, kpi_w, kpi_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, x, kpi_y, kpi_w, Inches(0.05), TEAL)
        tb_n = _tb(slide, x, kpi_y + Inches(0.06), kpi_w, Inches(0.55))
        _run(tb_n.text_frame, val, 30, bold=True, color=GOLD,
             font="Georgia", align=PP_ALIGN.CENTER, first=True)
        tb_l = _tb(slide, x, kpi_y + Inches(0.62), kpi_w, Inches(0.24))
        _run(tb_l.text_frame, lbl, 11, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, first=True)

    # Channel breakdown chart (left half)
    add_channel_chart(slide, overview.get("channel_breakdown", []),
                      Inches(0.30), Inches(2.05), Inches(4.40), Inches(3.20))

    # Goal breakdown table (right half)
    goals = overview.get("goal_breakdown", [])
    tx = Inches(4.95)
    ty = Inches(2.05)
    col_ws = [Inches(2.50), Inches(0.85), Inches(0.85), Inches(0.55)]
    headers = ["Channel / Goal", "Camps", "Segs", "Ops"]
    row_h = Inches(0.30)
    cx = tx
    for hdr, cw in zip(headers, col_ws):
        _rect(slide, cx, ty, cw, row_h, TEAL)
        tb = _tb(slide, cx + Inches(0.05), ty + Inches(0.04), cw - Inches(0.08), row_h - Inches(0.04))
        _run(tb.text_frame, hdr, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
        cx += cw
    for ri, g in enumerate(goals[:8]):
        ry = ty + (ri + 1) * row_h
        alt = KPI_BOX_BG if ri % 2 == 0 else RGBColor(0x1A, 0x3C, 0x56)
        vals = [g.get("goal", ""), str(g.get("count", "")), "", ""]
        cx = tx
        for i, (v, cw) in enumerate(zip(vals, col_ws)):
            _rect(slide, cx, ry, cw, row_h, alt, KPI_BOX_BD, 0.3)
            tb = _tb(slide, cx + Inches(0.05), ry + Inches(0.04), cw - Inches(0.08), row_h - Inches(0.04))
            txt_color = LIGHT_BLUE if i == 0 else WHITE
            _run(tb.text_frame, v, 9, color=txt_color, first=True)
            cx += cw


# ─── Slide 5: Condition Analysis ─────────────────────────────────────────────

def add_condition_analysis_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "03  CONDITION ANALYSIS — What Drives Audience Selection")

    analysis = report_data["condition_analysis"]
    conditions = analysis["top_conditions"]
    mix = analysis["condition_type_mix"]

    if conditions:
        chart_data = CategoryChartData()
        chart_data.categories = [c["condition_name"][:30] for c in conditions[:10]]
        chart_data.add_series("Usage", [c["usage_count"] for c in conditions[:10]])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.30), Inches(0.82), Inches(6.30), Inches(4.55),
            chart_data,
        ).chart
        chart.has_legend = False
        try:
            chart.plots[0].series[0].format.fill.solid()
            chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
        except Exception:
            pass

    # Right annotation boxes
    right_x = Inches(6.75)
    cdp = mix.get("CDP", 0); bq = mix.get("BQ", 0)
    box_data = [
        ("Standard Exclude",
         f"{conditions[0]['condition_name']}: {conditions[0]['usage_count']} uses" if conditions else "—",
         "Top #1 — used across most segments as standard exclusion."),
        ("Core Audience Pool",
         f"{conditions[1]['condition_name']}: {conditions[1]['usage_count']} uses" if len(conditions) > 1 else "—",
         "Base audience pool for this BU."),
        ("CDP vs BQ Balance",
         f"CDP: {cdp} uses  |  BQ: {bq} uses",
         "Balance indicates targeting sophistication."),
    ]
    for i, (title, val_line, desc_line) in enumerate(box_data):
        by = Inches(0.82) + i * Inches(1.20)
        _rect(slide, right_x, by, Inches(2.95), Inches(1.08), KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, right_x, by, Inches(2.95), Inches(0.28), TEAL)
        tb_t = _tb(slide, right_x + Inches(0.10), by + Inches(0.03), Inches(2.75), Inches(0.22))
        _run(tb_t.text_frame, title, 10, bold=True, color=WHITE, first=True)
        tb_v = _tb(slide, right_x + Inches(0.10), by + Inches(0.30), Inches(2.75), Inches(0.30))
        _run(tb_v.text_frame, val_line[:70], 8.5, color=GOLD, first=True)
        tb_d = _tb(slide, right_x + Inches(0.10), by + Inches(0.60), Inches(2.75), Inches(0.36))
        tb_d.text_frame.word_wrap = True
        _run(tb_d.text_frame, desc_line[:100], 8, color=LIGHT_BLUE, first=True)


# ─── Slide 6: Size Pool ───────────────────────────────────────────────────────

def add_size_pool_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "04  SIZE POOL DISTRIBUTION — Audience Size by Goal")

    pools = report_data.get("size_pool", [])
    card_w = Inches(4.50)
    card_h = Inches(0.85)
    gap = Inches(0.12)
    col_xs = [Inches(0.30), Inches(5.05)]
    start_y = Inches(0.82)

    for i, line in enumerate(pools[:8]):
        cx = col_xs[i % 2]
        cy = start_y + (i // 2) * (card_h + gap)
        lbl, val = (str(line).split(" median =", 1) + [""])[:2] if " median =" in str(line) else (str(line), "")
        _rect(slide, cx, cy, card_w, card_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, cx, cy, Inches(0.07), card_h, TEAL)
        tb_l = _tb(slide, cx + Inches(0.13), cy + Inches(0.05), card_w - Inches(0.15), Inches(0.25))
        _run(tb_l.text_frame, lbl.strip()[:60], 10, bold=True, color=GOLD, first=True)
        if val:
            tb_v = _tb(slide, cx + Inches(0.13), cy + Inches(0.33), card_w - Inches(0.15), Inches(0.45))
            tb_v.text_frame.word_wrap = True
            _run(tb_v.text_frame, f"Median = {val.strip()}", 8.5, color=LIGHT_BLUE, first=True)

    if not pools:
        tb = _tb(slide, Inches(0.30), Inches(1.00), Inches(9.40), Inches(4.50))
        _run(tb.text_frame, "No size pool data available.", 12, color=MUTED, first=True)


# ─── Slide 7: Operator Profiles ──────────────────────────────────────────────

def add_operator_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "05  PER-PERSON TACTICS — Operating Profiles")

    profiles = report_data.get("operator_profiles", [])
    accent_colors = [TEAL, GOLD, RGBColor(0x5E, 0xA8, 0xC5)]
    profile_labels = ["PROFILE A", "PROFILE B", "PROFILE C"]
    card_w = Inches(2.90)
    card_h = Inches(4.50)
    gap = Inches(0.25)
    start_y = Inches(0.85)

    for i, profile in enumerate(profiles[:3]):
        x = Inches(0.30) + i * (card_w + gap)
        accent = accent_colors[i % len(accent_colors)]
        _rect(slide, x, start_y, card_w, card_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, x, start_y, card_w, Inches(0.07), accent)
        tb_lbl = _tb(slide, x + Inches(0.10), start_y + Inches(0.12), card_w - Inches(0.20), Inches(0.25))
        _run(tb_lbl.text_frame, profile_labels[i], 8, bold=True, color=accent, first=True)
        tb_c = _tb(slide, x + Inches(0.10), start_y + Inches(0.42), card_w - Inches(0.20), Inches(3.90))
        tb_c.text_frame.word_wrap = True
        parts = str(profile).split("|") if "|" in str(profile) else [str(profile)]
        for j, part in enumerate(parts[:8]):
            _run(tb_c.text_frame, part.strip(), 9, color=LIGHT_BLUE, first=(j == 0))

    if not profiles:
        tb = _tb(slide, Inches(0.30), Inches(1.00), Inches(9.40), Inches(4.50))
        _run(tb.text_frame, "No operator profile data available.", 12, color=MUTED, first=True)


# ─── Slide 8: Quality Issues ─────────────────────────────────────────────────

def add_quality_issues_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "06  SEGMENT QUALITY ISSUES — Reuse Patterns & Data Gaps")

    issues = report_data.get("quality_issues", [])
    card_w = Inches(4.50)
    card_h = Inches(0.80)
    gap = Inches(0.12)
    col_xs = [Inches(0.30), Inches(5.05)]
    start_y = Inches(0.82)

    for i, issue in enumerate(issues[:8]):
        cx = col_xs[i % 2]
        cy = start_y + (i // 2) * (card_h + gap)
        parts = str(issue).split(".", 1)
        title = parts[0].strip()[:60]
        detail = parts[1].strip()[:120] if len(parts) > 1 else ""
        _rect(slide, cx, cy, card_w, card_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, cx, cy, Inches(0.07), card_h, GOLD)
        tb_t = _tb(slide, cx + Inches(0.13), cy + Inches(0.05), card_w - Inches(0.15), Inches(0.25))
        _run(tb_t.text_frame, title, 10, bold=True, color=GOLD, first=True)
        if detail:
            tb_d = _tb(slide, cx + Inches(0.13), cy + Inches(0.32), card_w - Inches(0.15), Inches(0.42))
            tb_d.text_frame.word_wrap = True
            _run(tb_d.text_frame, detail, 8.5, color=LIGHT_BLUE, first=True)

    if not issues:
        tb = _tb(slide, Inches(0.30), Inches(1.00), Inches(9.40), Inches(4.50))
        _run(tb.text_frame, "No quality issues detected.", 12, color=MUTED, first=True)


# ─── Slide 9: Recommendations ────────────────────────────────────────────────

def add_recommendations_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _inner_chrome(slide, "07  RECOMMENDATIONS — Actionable Improvements")

    recs = report_data.get("recommendations", [])
    priority_colors = [RED_HIGH, GOLD, TEAL, TEAL, TEAL, TEAL]
    priority_labels = ["HIGH", "HIGH", "MED", "MED", "LOW", "LOW"]
    card_h = Inches(0.72)
    gap = Inches(0.10)
    start_y = Inches(0.82)

    for i, rec in enumerate(recs[:6]):
        y = start_y + i * (card_h + gap)
        pc = priority_colors[i]
        pl = priority_labels[i]
        _rect(slide, Inches(0.30), y, Inches(9.40), card_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, Inches(0.30), y, Inches(0.07), card_h, pc)
        tb_n = _tb(slide, Inches(0.42), y + Inches(0.08), Inches(0.55), Inches(0.55))
        _run(tb_n.text_frame, f"R{i+1}", 20, bold=True, color=pc, font="Georgia", first=True)
        _rect(slide, Inches(1.05), y + Inches(0.22), Inches(0.55), Inches(0.26), pc)
        tb_p = _tb(slide, Inches(1.05), y + Inches(0.22), Inches(0.55), Inches(0.26))
        _run(tb_p.text_frame, pl, 7, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
        tb_r = _tb(slide, Inches(1.68), y + Inches(0.12), Inches(7.90), Inches(0.50))
        tb_r.text_frame.word_wrap = True
        _run(tb_r.text_frame, str(rec)[:200], 10, color=LIGHT_BLUE, first=True)

    if not recs:
        tb = _tb(slide, Inches(0.30), Inches(1.00), Inches(9.40), Inches(4.50))
        _run(tb.text_frame, "No recommendations available.", 12, color=MUTED, first=True)


# ─── Slide 10: Key Takeaways ─────────────────────────────────────────────────

def add_key_takeaways_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _rect(slide, Inches(0), Inches(0), Inches(0.14), SLIDE_H, TEAL)

    tb_t = _tb(slide, Inches(0.45), Inches(0.10), Inches(9.0), Inches(0.55))
    _run(tb_t.text_frame, "KEY TAKEAWAYS", 22, bold=True, color=WHITE, first=True)

    tb_sub = _tb(slide, Inches(0.45), Inches(0.65), Inches(9.0), Inches(0.32))
    _run(tb_sub.text_frame, f"{report_data['bu_name']} Team — Segment Analysis", 12,
         color=LIGHT_BLUE, first=True)

    _rect(slide, Inches(0.45), Inches(1.02), Inches(2.80), Inches(0.04), GOLD)

    takeaways = report_data.get("key_takeaways", [])
    item_y = Inches(1.15)
    item_h = Inches(0.72)
    gap = Inches(0.08)
    two_col = len(takeaways) > 3
    col_w = Inches(4.40) if two_col else Inches(9.10)

    for i, item in enumerate(takeaways[:6]):
        if two_col and i >= 3:
            x = Inches(5.20)
            y = item_y + (i - 3) * (item_h + gap)
        else:
            x = Inches(0.45)
            y = item_y + i * (item_h + gap)
        _rect(slide, x, y, col_w, item_h, KPI_BOX_BG, KPI_BOX_BD, 0.5)
        _rect(slide, x, y, Inches(0.07), item_h, TEAL)
        tb_n = _tb(slide, x + Inches(0.12), y + Inches(0.08), Inches(0.40), item_h - Inches(0.16))
        _run(tb_n.text_frame, f"{i+1:02d}", 16, bold=True, color=GOLD,
             font="Georgia", align=PP_ALIGN.CENTER, first=True)
        tb_item = _tb(slide, x + Inches(0.60), y + Inches(0.08), col_w - Inches(0.70), item_h - Inches(0.16))
        tb_item.text_frame.word_wrap = True
        _run(tb_item.text_frame, str(item)[:180], 9, color=LIGHT_BLUE, first=True)

    if not takeaways:
        tb = _tb(slide, Inches(0.45), Inches(1.15), Inches(9.10), Inches(4.20))
        _run(tb.text_frame, "No key takeaways available.", 12, color=MUTED, first=True)


# ─── Chart helper ─────────────────────────────────────────────────────────────

def add_channel_chart(slide, channel_breakdown, left, top, width, height):
    if not channel_breakdown:
        return
    chart_data = CategoryChartData()
    chart_data.categories = [x["channel"] for x in channel_breakdown[:6]]
    chart_data.add_series("Campaigns", [x["count"] for x in channel_breakdown[:6]])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    try:
        chart.plots[0].series[0].format.fill.solid()
        chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
    except Exception:
        pass


# ─── Entry point ──────────────────────────────────────────────────────────────

def write_ppt(report_data: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover_slide(prs, report_data)
    add_agenda_slide(prs, report_data)
    add_methodology_slide(prs, report_data)
    add_team_overview_slide(prs, report_data)
    add_condition_analysis_slide(prs, report_data)
    add_size_pool_slide(prs, report_data)
    add_operator_slide(prs, report_data)
    add_quality_issues_slide(prs, report_data)
    add_recommendations_slide(prs, report_data)
    add_key_takeaways_slide(prs, report_data)

    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    return str(output_file)
